from __future__ import annotations

import argparse
import json
import os
import re
import shutil
import sys
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Iterable


SUPPORTED_EXTENSIONS = {".docx", ".pptx", ".pdf", ".txt", ".md"}

# A page/slide/paragraph with fewer chars than this is treated as "almost no
# extractable text" -> a strong hint the content lives in images/equations.
LOW_TEXT_CHARS = 40
# Below this fraction of text-bearing pages, a PDF is flagged as likely scanned.
SCANNED_TEXT_RATIO = 0.5
# Keywords that mark a question as formula/derivation-heavy. If the KB has zero
# extracted formulas, answers to these questions cannot be courseware-grounded.
FORMULA_QUESTION_HINTS = (
    "公式", "导出", "推导", "表示式", "定义式", "证明", "计算", "方程",
    "formula", "derive", "equation", "expression",
)


@dataclass
class SourceBlock:
    source: str
    locator: str
    text: str


@dataclass
class MaterialDiagnostic:
    """Per-file extraction health, so 'image-heavy / scanned' is never silent."""
    source: str
    kind: str  # pdf | pptx | docx | txt | md
    units_total: int = 0          # pages / slides / paragraphs seen
    units_with_text: int = 0      # of those, how many yielded usable text
    chars_extracted: int = 0
    text_ratio: float = 0.0       # units_with_text / units_total
    likely_scanned: bool = False
    note: str = ""


@dataclass
class PageEvidence:
    """Visual evidence for image-heavy PDFs. Text blocks alone are not enough."""
    source: str
    page_number: int
    rendered_image_path: str
    text_chars: int = 0
    text_density: float = 0.0
    auto_keywords: list[str] = field(default_factory=list)
    manual_tags: list[str] = field(default_factory=list)


@dataclass
class KnowledgeBase:
    version: int
    materials: list[str]
    blocks: list[SourceBlock]
    formulas: list[str]
    format_rules: list[str]
    diagnostics: list[MaterialDiagnostic] = field(default_factory=list)
    page_evidence: list[PageEvidence] = field(default_factory=list)


def normalize_text(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _finalize_diag(diag: MaterialDiagnostic) -> MaterialDiagnostic:
    diag.text_ratio = (diag.units_with_text / diag.units_total) if diag.units_total else 0.0
    # A document is "likely scanned / image-heavy" when most of its units carry
    # almost no extractable text. This is the signal a teacher needs up front.
    diag.likely_scanned = (
        diag.units_total > 0 and diag.text_ratio < SCANNED_TEXT_RATIO
    )
    if diag.likely_scanned and not diag.note:
        diag.note = (
            f"only {diag.units_with_text}/{diag.units_total} units carry text "
            f"({diag.chars_extracted} chars) -> content likely lives in images/"
            f"equations; text-only grounding will be unreliable."
        )
    return diag


def extract_docx(path: Path) -> tuple[list[SourceBlock], MaterialDiagnostic]:
    from docx import Document

    doc = Document(str(path))
    blocks: list[SourceBlock] = []
    diag = MaterialDiagnostic(source=str(path), kind="docx")
    for idx, para in enumerate(doc.paragraphs, 1):
        diag.units_total += 1
        text = normalize_text(para.text)
        if text:
            diag.units_with_text += 1
            diag.chars_extracted += len(text)
            blocks.append(SourceBlock(str(path), f"paragraph {idx}", text))

    for table_idx, table in enumerate(doc.tables, 1):
        rows = []
        for row in table.rows:
            cells = [normalize_text(cell.text) for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            joined = "\n".join(rows)
            diag.chars_extracted += len(joined)
            blocks.append(SourceBlock(str(path), f"table {table_idx}", joined))
    return blocks, _finalize_diag(diag)


def extract_pptx(path: Path) -> tuple[list[SourceBlock], MaterialDiagnostic]:
    from pptx import Presentation

    prs = Presentation(str(path))
    blocks: list[SourceBlock] = []
    diag = MaterialDiagnostic(source=str(path), kind="pptx")
    for slide_idx, slide in enumerate(prs.slides, 1):
        diag.units_total += 1
        texts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = normalize_text(shape.text)
                if text:
                    texts.append(text)
            if getattr(shape, "has_table", False):
                table_rows = []
                for row in shape.table.rows:
                    cells = [normalize_text(cell.text) for cell in row.cells]
                    if any(cells):
                        table_rows.append(" | ".join(cells))
                if table_rows:
                    texts.append("\n".join(table_rows))
        text = normalize_text("\n".join(texts))
        if text:
            diag.units_with_text += 1
            diag.chars_extracted += len(text)
            blocks.append(SourceBlock(str(path), f"slide {slide_idx}", text))
    return blocks, _finalize_diag(diag)


def extract_pdf(path: Path) -> tuple[list[SourceBlock], MaterialDiagnostic]:
    diag = MaterialDiagnostic(source=str(path), kind="pdf")

    def harvest(pages_iter, page_count):
        diag.units_total = page_count
        for idx, get_text in pages_iter:
            text = normalize_text(get_text() or "")
            if text:
                diag.units_with_text += 1
                diag.chars_extracted += len(text)
                blocks.append(SourceBlock(str(path), f"page {idx}", text))

    blocks: list[SourceBlock] = []
    try:
        import pdfplumber

        with pdfplumber.open(str(path)) as pdf:
            pages = list(pdf.pages)
            harvest(
                ((i, (lambda p=p: p.extract_text())) for i, p in enumerate(pages, 1)),
                len(pages),
            )
        diag.note = "engine=pdfplumber"
    except Exception as exc:  # fall back, but record WHY
        from PyPDF2 import PdfReader

        reader = PdfReader(str(path))
        pages = list(reader.pages)
        diag.units_total = 0
        diag.units_with_text = 0
        diag.chars_extracted = 0
        blocks.clear()
        harvest(
            ((i, (lambda p=p: p.extract_text())) for i, p in enumerate(pages, 1)),
            len(pages),
        )
        diag.note = f"engine=PyPDF2 (pdfplumber failed: {type(exc).__name__})"
    return blocks, _finalize_diag(diag)


def pdf_page_texts(path: Path) -> tuple[list[str], str]:
    """Return per-page text. Empty strings are meaningful: likely scanned pages."""
    try:
        import pdfplumber

        with pdfplumber.open(str(path)) as pdf:
            return [normalize_text(page.extract_text() or "") for page in pdf.pages], "pdfplumber"
    except Exception as exc:
        from PyPDF2 import PdfReader

        reader = PdfReader(str(path))
        return [normalize_text(page.extract_text() or "") for page in reader.pages], f"PyPDF2 after pdfplumber {type(exc).__name__}"


def diagnose_pdf_file(path: Path) -> tuple[MaterialDiagnostic, list[int], str]:
    texts, engine = pdf_page_texts(path)
    diag = MaterialDiagnostic(source=str(path), kind="pdf", units_total=len(texts), note=f"engine={engine}")
    low_pages: list[int] = []
    for idx, text in enumerate(texts, 1):
        if text:
            diag.chars_extracted += len(text)
        if len(text) >= LOW_TEXT_CHARS:
            diag.units_with_text += 1
        else:
            low_pages.append(idx)
    return _finalize_diag(diag), low_pages, engine


def safe_pdf_stem(path: Path) -> str:
    stem = re.sub(r"[^A-Za-z0-9._-]+", "_", path.stem).strip("._-")
    return stem or "courseware"


def prepare_pdf_path(path: Path, out_dir: Path) -> Path:
    out_dir.mkdir(parents=True, exist_ok=True)
    dst = out_dir / f"{safe_pdf_stem(path)}.pdf"
    if path.resolve() != dst.resolve():
        shutil.copy2(path, dst)
    return dst


def render_pdf_pages(pdf_path: Path, out_dir: Path, pages: list[int], zoom: float = 2.0) -> list[PageEvidence]:
    import fitz

    out_dir.mkdir(parents=True, exist_ok=True)
    texts, _engine = pdf_page_texts(pdf_path)
    doc = fitz.open(pdf_path)
    evidence: list[PageEvidence] = []
    for page_no in pages:
        if page_no < 1 or page_no > len(doc):
            continue
        page = doc[page_no - 1]
        pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
        image_path = out_dir / f"page_{page_no:03}.png"
        pix.save(image_path)
        chars = len(texts[page_no - 1]) if page_no <= len(texts) else 0
        area = float(page.rect.width * page.rect.height) or 1.0
        evidence.append(
            PageEvidence(
                source=str(pdf_path),
                page_number=page_no,
                rendered_image_path=str(image_path),
                text_chars=chars,
                text_density=chars / area,
                auto_keywords=keyword_hits(texts[page_no - 1] if page_no <= len(texts) else ""),
            )
        )
    return evidence


def make_contact_sheets(pdf_path: Path, out_dir: Path, pages_per_sheet: int = 12) -> list[Path]:
    import fitz
    from PIL import Image, ImageDraw

    out_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(pdf_path)
    cols = 3
    rows = max(1, (pages_per_sheet + cols - 1) // cols)
    thumb_w = 360
    thumb_h = 270
    paths: list[Path] = []
    for start in range(0, len(doc), pages_per_sheet):
        page_indexes = list(range(start, min(start + pages_per_sheet, len(doc))))
        sheet = Image.new("RGB", (cols * thumb_w, rows * (thumb_h + 28)), "white")
        draw = ImageDraw.Draw(sheet)
        for n, idx in enumerate(page_indexes):
            page = doc[idx]
            pix = page.get_pixmap(matrix=fitz.Matrix(0.38, 0.38), alpha=False)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            img.thumbnail((thumb_w, thumb_h), Image.LANCZOS)
            x = (n % cols) * thumb_w + (thumb_w - img.width) // 2
            y = (n // cols) * (thumb_h + 28)
            sheet.paste(img, (x, y))
            draw.text(((n % cols) * thumb_w + 8, y + thumb_h + 4), f"Page {idx + 1}", fill="black")
        path = out_dir / f"contact_{start + 1:03}_{page_indexes[-1] + 1:03}.jpg"
        sheet.save(path, quality=88)
        paths.append(path)
    return paths


KEYWORD_ALIASES = {
    "stokes": ("stokes", "斯托克斯", "偏振"),
    "polarization": ("polarization", "偏振", "线偏振", "圆偏振"),
    "retarded_potential": ("推迟势", "retarded potential", "洛伦兹规范"),
    "angular_distribution": ("角分布", "辐射角", "beaming", "相对论"),
    "formula": ("公式", "导出", "推导", "表示式", "方程"),
}


def keyword_hits(text: str) -> list[str]:
    low = text.lower()
    hits = []
    for label, aliases in KEYWORD_ALIASES.items():
        if any(alias.lower() in low for alias in aliases):
            hits.append(label)
    return hits


def locate_keyword_pages(pdf_path: Path, query: str, max_pages: int = 20) -> list[tuple[int, int, list[str], str]]:
    texts, _engine = pdf_page_texts(pdf_path)
    q_tokens = token_set(query)
    results: list[tuple[int, int, list[str], str]] = []
    for idx, text in enumerate(texts, 1):
        hits = keyword_hits(text)
        score = len(q_tokens & token_set(text)) + len(hits) * 2
        if score > 0:
            snippet = text[:220].replace("\n", " ")
            results.append((idx, score, hits, snippet))
    return sorted(results, key=lambda item: item[1], reverse=True)[:max_pages]


def run_ocr_on_images(image_paths: list[Path], lang: str = "chi_sim+eng") -> tuple[dict[str, str], str]:
    """Best-effort local OCR. If no engine exists, return an honest error string."""
    if not shutil.which("tesseract"):
        return {}, "OCR unavailable: tesseract executable not found. Install Tesseract or use manual page-image review."
    try:
        import pytesseract
        from PIL import Image
    except Exception as exc:
        return {}, f"OCR unavailable: missing Python OCR dependencies ({type(exc).__name__})."

    outputs: dict[str, str] = {}
    for image_path in image_paths:
        try:
            outputs[str(image_path)] = pytesseract.image_to_string(Image.open(image_path), lang=lang)
        except Exception as exc:
            outputs[str(image_path)] = f"[OCR failed: {type(exc).__name__}: {exc}]"
    return outputs, "ok"


def extract_plain(path: Path) -> tuple[list[SourceBlock], MaterialDiagnostic]:
    kind = path.suffix.lower().lstrip(".") or "txt"
    text = normalize_text(path.read_text(encoding="utf-8", errors="ignore"))
    diag = MaterialDiagnostic(source=str(path), kind=kind, units_total=1)
    blocks: list[SourceBlock] = []
    if text:
        diag.units_with_text = 1
        diag.chars_extracted = len(text)
        blocks.append(SourceBlock(str(path), "text", text))
    return blocks, _finalize_diag(diag)


def extract_material(path: Path) -> tuple[list[SourceBlock], MaterialDiagnostic]:
    ext = path.suffix.lower()
    if ext == ".docx":
        return extract_docx(path)
    if ext == ".pptx":
        return extract_pptx(path)
    if ext == ".pdf":
        return extract_pdf(path)
    if ext in {".txt", ".md"}:
        return extract_plain(path)
    raise ValueError(f"Unsupported material type: {path}")


def chunk_block(block: SourceBlock, max_chars: int = 1200) -> list[SourceBlock]:
    text = block.text
    if len(text) <= max_chars:
        return [block]

    parts = re.split(r"(?<=[。！？.!?])\s+|\n{2,}", text)
    chunks: list[SourceBlock] = []
    current = ""
    chunk_idx = 1
    for part in parts:
        part = part.strip()
        if not part:
            continue
        if current and len(current) + len(part) + 1 > max_chars:
            chunks.append(SourceBlock(block.source, f"{block.locator} chunk {chunk_idx}", current))
            chunk_idx += 1
            current = part
        else:
            current = f"{current}\n{part}".strip()
    if current:
        chunks.append(SourceBlock(block.source, f"{block.locator} chunk {chunk_idx}", current))
    return chunks


def find_formulas(blocks: Iterable[SourceBlock]) -> list[str]:
    candidates: list[str] = []
    formula_line = re.compile(r"(^|[\s:：])[\wα-ωΑ-ΩµπθλβγσρΔΩ]+.*[=≈∝≤≥<>].+")
    for block in blocks:
        for line in block.text.splitlines():
            line = normalize_text(line)
            if 4 <= len(line) <= 240 and formula_line.search(line):
                candidates.append(line)
    return unique_preserve_order(candidates)


def find_format_rules(blocks: Iterable[SourceBlock]) -> list[str]:
    keywords = ("格式", "步骤", "要求", "作答", "答案", "解法", "变量", "记为", "必须", "应当", "注意")
    rules = []
    for block in blocks:
        for line in block.text.splitlines():
            line = normalize_text(line)
            if 6 <= len(line) <= 260 and any(keyword in line for keyword in keywords):
                rules.append(line)
    return unique_preserve_order(rules)


def unique_preserve_order(items: Iterable[str]) -> list[str]:
    seen = set()
    result = []
    for item in items:
        key = item.strip()
        if key and key not in seen:
            result.append(key)
            seen.add(key)
    return result


def token_set(text: str) -> set[str]:
    lowered = text.lower()
    words = set(re.findall(r"[a-z0-9_]+", lowered))
    cjk = re.findall(r"[\u4e00-\u9fff]", lowered)
    grams = {"".join(cjk[i : i + 2]) for i in range(max(0, len(cjk) - 1))}
    return words | grams


def score_match(question: str, block: SourceBlock) -> float:
    q_tokens = token_set(question)
    b_tokens = token_set(block.text)
    if not q_tokens or not b_tokens:
        return 0.0
    overlap = len(q_tokens & b_tokens)
    return overlap / (len(q_tokens) ** 0.5)


def retrieve(blocks: list[SourceBlock], question: str, top_k: int) -> list[SourceBlock]:
    ranked = sorted(((score_match(question, block), block) for block in blocks), key=lambda x: x[0], reverse=True)
    return [block for score, block in ranked[:top_k] if score > 0]


def collect_files(materials: list[Path]) -> list[Path]:
    files: list[Path] = []
    for material in materials:
        if material.is_dir():
            files.extend(
                p for p in material.rglob("*") if p.suffix.lower() in SUPPORTED_EXTENSIONS
            )
        else:
            files.append(material)
    return files


def build_kb(materials: list[Path], render_page_evidence: bool = False, page_evidence_dir: Path | None = None) -> KnowledgeBase:
    raw_blocks: list[SourceBlock] = []
    diagnostics: list[MaterialDiagnostic] = []
    page_evidence: list[PageEvidence] = []
    for file in collect_files(materials):
        blocks, diag = extract_material(file)
        raw_blocks.extend(blocks)
        diagnostics.append(diag)
        if render_page_evidence and file.suffix.lower() == ".pdf":
            target_dir = page_evidence_dir or Path("standard_answer_harness/out/page_evidence") / safe_pdf_stem(file)
            try:
                diag_pdf, _low_pages, _engine = diagnose_pdf_file(file)
                pages = list(range(1, diag_pdf.units_total + 1))
                page_evidence.extend(render_pdf_pages(file, target_dir, pages, zoom=1.3))
            except Exception as exc:
                page_evidence.append(
                    PageEvidence(
                        source=str(file),
                        page_number=0,
                        rendered_image_path="",
                        manual_tags=[f"page evidence render failed: {type(exc).__name__}: {exc}"],
                    )
                )

    blocks = [chunk for block in raw_blocks for chunk in chunk_block(block)]
    return KnowledgeBase(
        version=1,
        materials=[str(p) for p in materials],
        blocks=blocks,
        formulas=find_formulas(blocks),
        format_rules=find_format_rules(blocks),
        diagnostics=diagnostics,
        page_evidence=page_evidence,
    )


def save_kb(kb: KnowledgeBase, path: Path) -> None:
    payload = asdict(kb)
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def load_kb(path: Path) -> KnowledgeBase:
    payload = json.loads(path.read_text(encoding="utf-8"))
    payload["blocks"] = [SourceBlock(**block) for block in payload["blocks"]]
    # Backward compatible: KBs built before diagnostics existed simply have none.
    payload["diagnostics"] = [
        MaterialDiagnostic(**d) for d in payload.get("diagnostics", [])
    ]
    payload["page_evidence"] = [
        PageEvidence(**p) for p in payload.get("page_evidence", [])
    ]
    return KnowledgeBase(**payload)


def read_questions(path: Path) -> list[str]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    parts = re.split(r"\n\s*\n|^\s*\d+[.、]\s*", text, flags=re.M)
    return [normalize_text(part) for part in parts if normalize_text(part)]


# Minimum total evidence chars below which retrieval is treated as too thin to
# carry a real answer (the matched blocks are fragments like "谱分布：").
MIN_EVIDENCE_CHARS = 120
# Honest self-labels an answer can use to declare it is NOT courseware-grounded.
# If an answer is ungrounded but contains none of these, the harness blocks it.
HONESTY_MARKERS = (
    "课件证据不足", "未接地", "证据不足", "模型先验", "模型推导",
    "无法据课件", "课件中未", "未能从课件", "可能与资料不一致",
    "可能和资料不一致", "非课件直接依据", "与课件不一致风险",
)


@dataclass
class GroundingVerdict:
    """Whether the courseware actually supports answering this question."""
    level: str  # grounded | weak | ungrounded
    is_formula_question: bool
    evidence_count: int
    evidence_chars: int
    kb_formulas: int
    reasons: list[str] = field(default_factory=list)


def is_formula_question(question: str) -> bool:
    low = question.lower()
    return any(h.lower() in low for h in FORMULA_QUESTION_HINTS)


def assess_grounding(question: str, evidence: list[SourceBlock], kb: KnowledgeBase) -> GroundingVerdict:
    """The honesty core: decide if the KB can actually ground this answer.

    This is intentionally conservative. A 'standard answer' tool must not imply
    courseware backing it does not have, especially for image-heavy/scanned decks
    where the physics lives in pictures and the text layer is near-empty.
    """
    formula_q = is_formula_question(question)
    ev_chars = sum(len(b.text) for b in evidence)
    reasons: list[str] = []

    if not evidence:
        reasons.append("未检索到任何课件证据块。")
        return GroundingVerdict("ungrounded", formula_q, 0, 0, len(kb.formulas), reasons)

    # A formula/derivation question cannot be courseware-grounded if the KB has
    # zero extracted formulas: the answer's equations are necessarily the model's
    # own prior, not the courseware's.
    if formula_q and not kb.formulas:
        reasons.append(
            "公式/推导类题目，但知识库自动抽取到的公式数为 0 "
            "（课件公式很可能在图片/公式对象里，未进入文本层）。"
            "本题的公式无法由课件接地。"
        )
        return GroundingVerdict("ungrounded", True, len(evidence), ev_chars, 0, reasons)

    if ev_chars < MIN_EVIDENCE_CHARS:
        reasons.append(
            f"匹配证据过薄（合计约 {ev_chars} 字，低于 {MIN_EVIDENCE_CHARS} 字阈值），"
            "多为标题/碎片，不足以支撑完整作答。"
        )
        return GroundingVerdict("weak", formula_q, len(evidence), ev_chars, len(kb.formulas), reasons)

    if formula_q and kb.formulas:
        reasons.append("公式类题目，且知识库含可引用公式，可尝试课件接地。")
    else:
        reasons.append("证据量足以支撑概念性作答。")
    return GroundingVerdict("grounded", formula_q, len(evidence), ev_chars, len(kb.formulas), reasons)


def provenance_banner(verdict: GroundingVerdict) -> str:
    """A required, machine-checkable header so the answer's grounding is explicit."""
    label = {
        "grounded": "课件接地（自动抽取证据足够）",
        "weak": "弱接地（仅碎片证据，需人工/页图核对）",
        "ungrounded": "未接地（仅模型先验，课件未提供可引用依据）",
    }[verdict.level]
    lines = [
        "## 接地度诊断",
        f"- 接地状态：**{label}**",
        f"- 命中证据块：{verdict.evidence_count}　证据字数：{verdict.evidence_chars}　"
        f"知识库公式数：{verdict.kb_formulas}",
    ]
    for r in verdict.reasons:
        lines.append(f"- {r}")
    if verdict.level != "grounded":
        lines.append(
            "- **出处标注要求**：本题答案中的公式/结论请按 "
            "`[自动抽取] / [页图核对] / [外部资料] / [模型推导]` 标注来源，"
            "不得默认呈现为课件原文。"
        )
    return "\n".join(lines)


def kb_has_scanned_material(kb: KnowledgeBase) -> bool:
    return any(d.likely_scanned for d in kb.diagnostics)


def formula_extraction_warning(kb: KnowledgeBase, verdict: GroundingVerdict | None = None) -> str:
    if kb.formulas:
        return ""
    scanned = kb_has_scanned_material(kb)
    if not scanned and (verdict is None or not verdict.is_formula_question):
        return ""
    return (
        "## 公式识别状态\n"
        "- **未识别到任何课件公式。**\n"
        "- 当前系统没有从文本层/OCR 中抽取出公式；如果课件是扫描页、图片公式或 Office 公式，"
        "这里不能证明后续公式来自课件。\n"
        "- 因此本题不能标为“课件标准答案”。如出现公式，只能标注为 `[模型推导]`、`[外部资料]` "
        "或 `[页图核对]`，并需要人工/视觉模型核对。"
    )


def _grounding_instructions(verdict: GroundingVerdict | None) -> str:
    """Inject grounding-state-specific rules into the prompt so the model is
    forced to label provenance honestly instead of dressing up its own prior
    as courseware."""
    if verdict is None:
        return ""
    if verdict.level == "ungrounded":
        return (
            "\n接地状态：未接地（课件未提供可引用依据）。\n"
            "硬性要求：\n"
            "- 允许给出 AI 自己的推导、判断或补充，但必须明确声明："
            "“以下内容可能与资料不一致，非课件直接依据”。\n"
            "- 凡课件未直接给出的公式/结论，必须逐条标注 `[模型推导]` 或 `[外部资料]`，"
            "不得标注或暗示为课件原文。\n"
            "- 不得虚构课件页码、公式编号或原文表述。\n"
        )
    if verdict.level == "weak":
        return (
            "\n接地状态：弱接地（仅碎片证据）。\n"
            "硬性要求：\n"
            "- 仅碎片/标题类证据可引用，超出部分必须标注 `[模型推导]` 或 `[页图核对]`。\n"
            "- 如需补全连贯推导，可以加入 AI 自己的想法，但必须提示："
            "“补全部分可能与资料不一致”。\n"
        )
    return (
        "\n接地状态：课件接地。可优先沿用课件证据中的公式与表述，"
        "并在“依据”部分列出证据编号。\n"
    )


def make_prompt(question: str, evidence: list[SourceBlock], kb: KnowledgeBase, external_context: str = "", verdict: GroundingVerdict | None = None) -> str:
    evidence_text = "\n\n".join(
        f"[{idx}] {Path(block.source).name} {block.locator}\n{block.text}" for idx, block in enumerate(evidence, 1)
    )
    formulas = "\n".join(f"- {formula}" for formula in kb.formulas[:50]) or "- 无自动抽取公式"
    rules = "\n".join(f"- {rule}" for rule in kb.format_rules[:50]) or "- 无自动抽取格式规则"
    external = f"\n\n可选拓展资料，只有在不违背课件时可使用：\n{external_context}" if external_context else ""
    grounding = _grounding_instructions(verdict)
    return f"""你是授课教师的标准答案生成助手。目标是尽量生成匹配课件口径的标准答案，但必须诚实区分课件依据和 AI 补全。

任务：
1. 先复述题目，标题必须为“题目复述”。
2. 优先沿用课件中的公式、变量名、术语、解题步骤和表述。
3. 允许加入 AI 自己的推导、判断或补充；但凡非课件直接给出的内容，必须说明“可能与资料不一致”。
4. 输出结构：题目复述、标准答案、依据、资料不一致风险、合规自检。
5. 依据部分列出使用的证据编号；不得虚构课件页码、公式编号或原文表述。
6. 凡非课件直接给出的内容，必须按 `[自动抽取] / [页图核对] / [外部资料] / [模型推导]` 标注来源。
{grounding}
课件中自动抽取的公式：
{formulas}

课件中自动抽取的作答/格式规则：
{rules}

课件证据：
{evidence_text}
{external}

题目：
{question}
"""


def call_openai(prompt: str, model: str) -> str:
    from openai import OpenAI

    client = OpenAI()
    response = client.responses.create(
        model=model,
        input=prompt,
    )
    return response.output_text


def offline_answer(question: str, evidence: list[SourceBlock], kb: KnowledgeBase, verdict: GroundingVerdict | None = None) -> str:
    formula_warning = formula_extraction_warning(kb, verdict)
    if not evidence:
        return f"""{formula_warning}

## 风险草稿
课件证据不足。离线模式不会自行补全答案；如需 AI 补全，请启用模型，并把补全部分标注为 `[模型推导]`，同时说明可能与资料不一致。

## 资料不一致风险
- 未找到匹配课件证据，任何完整答案都不是课件直接依据。

## 合规自检
- 未找到匹配课件证据。
"""

    cited = "\n".join(
        f"- [{idx}] {Path(block.source).name} {block.locator}: {block.text[:350]}"
        for idx, block in enumerate(evidence, 1)
    )
    rules = "\n".join(f"- {rule}" for rule in kb.format_rules[:8]) or "- 未抽取到显式格式规则。"
    formulas = "\n".join(f"- {formula}" for formula in kb.formulas[:8]) or "- 未抽取到显式公式。"
    # The offline path never invents physics. When grounding is weak/ungrounded it
    # says so plainly instead of producing a confident-looking draft.
    if verdict is not None and verdict.level != "grounded":
        honest = "课件证据不足" if verdict.level == "ungrounded" else "课件证据较薄，需人工/页图核对"
        return f"""{formula_warning}

## 风险草稿
{honest}。未调用大模型，仅列出可引用的课件检索证据；如需完整作答，可由 AI 补全，但必须标注 `[模型推导]` 并说明可能与资料不一致：

{cited}

## 课件公式候选
{formulas}

## 作答格式候选
{rules}

## 合规自检
- 接地状态：{verdict.level}。
- 本草稿仅含课件检索碎片，任何超出部分须标注 `[模型推导]`/`[外部资料]`，并提示可能与资料不一致。
"""
    return f"""## 标准答案草稿
当前未调用大模型，仅给出基于课件检索的教师用草稿。请依据下列证据组织最终答案：

{cited}

## 课件公式候选
{formulas}

## 作答格式候选
{rules}

## 合规自检
- 已限制在课件检索证据内。
- 需要人工或启用 OpenAI 模型补全连贯解题过程。
"""


@dataclass
class ComplianceFinding:
    severity: str  # block | warn | info
    message: str


def compliance_check(answer: str, evidence: list[SourceBlock], kb: KnowledgeBase, verdict: GroundingVerdict | None = None) -> list[ComplianceFinding]:
    """A blocking hard contract, not a rubber stamp.

    The gate allows AI supplementation. What it blocks is unlabeled
    supplementation that looks like courseware-grounded truth.
    """
    findings: list[ComplianceFinding] = []
    has_honest_marker = any(marker in answer for marker in HONESTY_MARKERS)

    if "题目复述" not in answer and "## 题目" not in answer:
        findings.append(ComplianceFinding("info", "模型正文未含题目复述；harness 外层会自动补上。"))

    if not evidence and not has_honest_marker:
        findings.append(ComplianceFinding("block", "没有匹配证据但未说明答案可能与资料不一致。"))

    if verdict is not None:
        if verdict.is_formula_question and kb_has_scanned_material(kb) and not kb.formulas:
            findings.append(ComplianceFinding(
                "block",
                "扫描/图片型课件未识别到任何公式；本题是公式/推导类，不能标为课件标准答案，只能作为风险草稿或经页图/OCR/人工核对后发布。",
            ))
        if verdict.level == "ungrounded" and not has_honest_marker:
            findings.append(ComplianceFinding(
                "block",
                "接地度诊断为“未接地”，但答案未声明 "
                f"{ '/'.join(HONESTY_MARKERS[:3]) } 等诚实标记——"
                "这会把 AI 补全冒充成课件标准答案。",
            ))
        elif verdict.level == "weak" and not has_honest_marker:
            findings.append(ComplianceFinding(
                "warn",
                "接地度诊断为“弱接地”，建议在答案中显式区分课件证据与模型补全部分。",
            ))
        if verdict.is_formula_question and not kb.formulas and not has_honest_marker:
            findings.append(ComplianceFinding(
                "block",
                "公式/推导类题目且知识库公式数为 0，答案却未标注公式来源——"
                "公式无法由课件接地，必须逐条标注 `[模型推导]`/`[外部资料]`，并说明可能与资料不一致。",
            ))

    for formula in kb.formulas[:50]:
        key = formula.split("=")[0].strip()
        if key and key in answer and formula not in answer:
            findings.append(ComplianceFinding(
                "warn",
                f"答案可能使用了公式变量“{key}”，但未完整沿用课件公式：{formula}",
            ))

    if evidence and not any(marker in answer for marker in ("依据", "证据", "[1]")):
        findings.append(ComplianceFinding("warn", "答案缺少课件依据引用。"))

    return findings


def command_build(args: argparse.Namespace) -> int:
    materials = [Path(p).resolve() for p in args.materials]
    missing = [str(p) for p in materials if not p.exists()]
    if missing:
        print(f"Missing materials: {', '.join(missing)}", file=sys.stderr)
        return 2
    kb = build_kb(
        materials,
        render_page_evidence=args.render_page_evidence,
        page_evidence_dir=Path(args.page_evidence_dir).resolve() if args.page_evidence_dir else None,
    )
    save_kb(kb, Path(args.out).resolve())
    print(f"Built KB: {args.out}")
    print(f"Blocks: {len(kb.blocks)}")
    print(f"Formulas: {len(kb.formulas)}")
    print(f"Format rules: {len(kb.format_rules)}")
    print(f"Page image evidence: {len(kb.page_evidence)}")
    print(f"Page image evidence: {len(kb.page_evidence)}")

    # Surface extraction health up front. A teacher must know BEFORE answering
    # that a deck is image-heavy/scanned, otherwise the KB is a confident-looking
    # empty shell.
    scanned = [d for d in kb.diagnostics if d.likely_scanned]
    print("\nExtraction diagnostics:")
    for d in kb.diagnostics:
        flag = "  [LIKELY SCANNED/IMAGE-HEAVY]" if d.likely_scanned else ""
        print(
            f"- {Path(d.source).name} ({d.kind}): "
            f"{d.units_with_text}/{d.units_total} units with text, "
            f"{d.chars_extracted} chars, ratio={d.text_ratio:.2f}{flag}"
        )
        if d.note:
            print(f"    note: {d.note}")
    if scanned:
        print(
            f"\nWARNING: {len(scanned)} material(s) look scanned/image-heavy. "
            "Text-only grounding will be unreliable; verify against page images "
            "and do not treat formula answers as courseware-grounded.",
            file=sys.stderr,
        )
    if not kb.formulas:
        print(
            "WARNING: zero formulas were auto-extracted. Formula/derivation "
            "questions cannot be courseware-grounded from this KB.",
            file=sys.stderr,
        )
    return 0


def command_diagnose_pdf(args: argparse.Namespace) -> int:
    pdf = Path(args.pdf).resolve()
    if not pdf.exists():
        print(f"Missing PDF: {pdf}", file=sys.stderr)
        return 2
    diag, low_pages, engine = diagnose_pdf_file(pdf)
    print(f"PDF: {pdf}")
    print(f"Pages: {diag.units_total}")
    print(f"Extraction engine: {engine}")
    print(f"Pages with usable text: {diag.units_with_text}/{diag.units_total}")
    print(f"Extracted chars: {diag.chars_extracted}")
    print(f"Text ratio: {diag.text_ratio:.2f}")
    print(f"Likely scanned/image-heavy: {diag.likely_scanned}")
    if not diag.likely_scanned:
        print("Verdict: copyable-text PDF. Text harness can be useful.")
    else:
        print(
            "Verdict: scanned/image-heavy PDF. Text harness is unreliable; "
            "use prep-pdf, page-image evidence, and OCR/manual visual review."
        )
    if low_pages:
        preview = ", ".join(str(p) for p in low_pages[:30])
        suffix = " ..." if len(low_pages) > 30 else ""
        print(f"Low/empty text pages: {preview}{suffix}")
    return 0


def command_prep_pdf(args: argparse.Namespace) -> int:
    pdf = Path(args.pdf).resolve()
    if not pdf.exists():
        print(f"Missing PDF: {pdf}", file=sys.stderr)
        return 2
    out_dir = Path(args.out).resolve()
    safe_pdf = prepare_pdf_path(pdf, out_dir)
    diag, low_pages, engine = diagnose_pdf_file(safe_pdf)

    last_count = max(1, args.last_pages)
    last_pages = list(range(max(1, diag.units_total - last_count + 1), diag.units_total + 1))
    rendered = render_pdf_pages(safe_pdf, out_dir / "rendered_last_pages", last_pages, zoom=args.zoom)
    contacts = make_contact_sheets(safe_pdf, out_dir / "contact", pages_per_sheet=args.pages_per_sheet)

    manifest = {
        "source_pdf": str(pdf),
        "safe_pdf": str(safe_pdf),
        "pages": diag.units_total,
        "engine": engine,
        "pages_with_text": diag.units_with_text,
        "chars_extracted": diag.chars_extracted,
        "text_ratio": diag.text_ratio,
        "likely_scanned": diag.likely_scanned,
        "last_page_images": [asdict(p) for p in rendered],
        "contact_sheets": [str(p) for p in contacts],
        "recommendation": (
            "Use page-image/OCR workflow; do not trust text-only grounding."
            if diag.likely_scanned else
            "Text extraction is usable; build KB is reasonable."
        ),
    }
    manifest_path = out_dir / "prep_manifest.json"
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")

    print(f"Prepared PDF workspace: {out_dir}")
    print(f"Safe PDF: {safe_pdf}")
    print(f"Pages: {diag.units_total}  Text ratio: {diag.text_ratio:.2f}  Chars: {diag.chars_extracted}")
    print(f"Likely scanned/image-heavy: {diag.likely_scanned}")
    print("Rendered last pages:")
    for ev in rendered:
        print(f"- page {ev.page_number}: {ev.rendered_image_path}")
    print(f"Contact sheets: {len(contacts)} written under {out_dir / 'contact'}")
    print(f"Manifest: {manifest_path}")
    if diag.likely_scanned:
        print("Next step: run homework-last-page or inspect contact sheets, then OCR only selected pages.")
    return 0


def command_homework_last_page(args: argparse.Namespace) -> int:
    pdf = Path(args.pdf).resolve()
    if not pdf.exists():
        print(f"Missing PDF: {pdf}", file=sys.stderr)
        return 2
    out_dir = Path(args.out).resolve()
    safe_pdf = prepare_pdf_path(pdf, out_dir)
    diag, _low_pages, _engine = diagnose_pdf_file(safe_pdf)
    pages = list(range(max(1, diag.units_total - args.last_pages + 1), diag.units_total + 1))
    rendered = render_pdf_pages(safe_pdf, out_dir / "homework_last_pages", pages, zoom=args.zoom)
    texts, _ = pdf_page_texts(safe_pdf)

    print(f"PDF pages: {diag.units_total}")
    print(f"Likely scanned/image-heavy: {diag.likely_scanned}")
    print("Last-page homework candidates:")
    for ev in rendered:
        text = texts[ev.page_number - 1] if ev.page_number <= len(texts) else ""
        print(f"\n--- page {ev.page_number} ---")
        print(f"image: {ev.rendered_image_path}")
        if text:
            print("auto text:")
            print(text[:1200])
        else:
            print("auto text: [none extracted]")

    if args.ocr:
        ocr_texts, status = run_ocr_on_images([Path(ev.rendered_image_path) for ev in rendered], lang=args.ocr_lang)
        print(f"\nOCR status: {status}")
        for path, text in ocr_texts.items():
            print(f"\n--- OCR {path} ---")
            print(text[:2000])
    elif diag.likely_scanned:
        print("\nOCR not run. For scanned PDFs, inspect the page image or rerun with --ocr if Tesseract is installed.")
    return 0


def parse_pages(spec: str) -> list[int]:
    pages: set[int] = set()
    for part in spec.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            start_s, end_s = part.split("-", 1)
            start, end = int(start_s), int(end_s)
            pages.update(range(min(start, end), max(start, end) + 1))
        else:
            pages.add(int(part))
    return sorted(pages)


def command_ocr_pages(args: argparse.Namespace) -> int:
    pdf = Path(args.pdf).resolve()
    if not pdf.exists():
        print(f"Missing PDF: {pdf}", file=sys.stderr)
        return 2
    out_dir = Path(args.out).resolve()
    safe_pdf = prepare_pdf_path(pdf, out_dir)
    pages = parse_pages(args.pages)
    rendered = render_pdf_pages(safe_pdf, out_dir / "ocr_pages", pages, zoom=args.zoom)
    ocr_texts, status = run_ocr_on_images([Path(ev.rendered_image_path) for ev in rendered], lang=args.ocr_lang)
    out_text = out_dir / "ocr_pages.txt"
    lines = [f"OCR status: {status}", ""]
    for ev in rendered:
        image_path = ev.rendered_image_path
        lines.append(f"--- page {ev.page_number} ---")
        lines.append(f"image: {image_path}")
        lines.append(ocr_texts.get(image_path, "[no OCR output]"))
        lines.append("")
    out_text.write_text("\n".join(lines), encoding="utf-8")
    print(f"Rendered pages: {[ev.page_number for ev in rendered]}")
    print(f"OCR status: {status}")
    print(f"OCR output: {out_text}")
    return 0


def command_locate_pages(args: argparse.Namespace) -> int:
    pdf = Path(args.pdf).resolve()
    if not pdf.exists():
        print(f"Missing PDF: {pdf}", file=sys.stderr)
        return 2
    matches = locate_keyword_pages(pdf, args.query, max_pages=args.max_pages)
    print(f"Query: {args.query}")
    if not matches:
        print("No text-layer keyword matches. This often means the PDF is scanned; use prep-pdf contact sheets or OCR selected pages.")
        return 1
    for page_no, score, hits, snippet in matches:
        print(f"- page {page_no}: score={score}, auto_keywords={hits or []}")
        if snippet:
            print(f"  {snippet}")
    return 0


def command_diagnose(args: argparse.Namespace) -> int:
    """Inspect a KB's extraction health without re-building it. Answers the only
    question that matters before trusting answers: can this courseware actually
    ground anything, or is it an empty shell?"""
    kb = load_kb(Path(args.kb).resolve())
    print(f"Materials: {len(kb.materials)}")
    print(f"Blocks: {len(kb.blocks)}　Formulas: {len(kb.formulas)}　Format rules: {len(kb.format_rules)}")

    if not kb.diagnostics:
        print(
            "\nNo per-file diagnostics in this KB (built by an older version). "
            "Re-run `build` to populate extraction health.",
            file=sys.stderr,
        )
    else:
        print("\nPer-material extraction health:")
        for d in kb.diagnostics:
            flag = "  [LIKELY SCANNED/IMAGE-HEAVY]" if d.likely_scanned else ""
            print(
                f"- {Path(d.source).name} ({d.kind}): "
                f"{d.units_with_text}/{d.units_total} units with text, "
                f"{d.chars_extracted} chars, ratio={d.text_ratio:.2f}{flag}"
            )
            if d.note:
                print(f"    note: {d.note}")

    # An overall verdict on whether the KB can ground a formula-heavy course.
    scanned = [d for d in kb.diagnostics if d.likely_scanned]
    print("\nOverall verdict:")
    if not kb.blocks:
        print("- UNUSABLE: no text blocks extracted at all.")
    elif not kb.formulas and scanned:
        print("- WEAK SHELL: text layer is sparse and no formulas extracted; "
              "treat as image-heavy. Formula answers cannot be courseware-grounded.")
    elif not kb.formulas:
        print("- PARTIAL: text extracted but zero formulas; concept questions may "
              "ground, formula/derivation questions cannot.")
    else:
        print("- OK: text and formulas present; grounding is plausible (still "
              "verify retrieval quality per question).")
    return 0


def command_answer(args: argparse.Namespace) -> int:
    kb = load_kb(Path(args.kb).resolve())
    questions = read_questions(Path(args.questions).resolve())
    external_context = Path(args.external_context).read_text(encoding="utf-8", errors="ignore") if args.external_context else ""
    model = args.model or os.getenv("STANDARD_ANSWER_MODEL")

    sections = []
    blocked_total = 0
    for idx, question in enumerate(questions, 1):
        evidence = retrieve(kb.blocks, question, args.top_k)
        verdict = assess_grounding(question, evidence, kb)
        if model:
            prompt = make_prompt(question, evidence, kb, external_context, verdict)
            answer = call_openai(prompt, model)
        else:
            answer = offline_answer(question, evidence, kb, verdict)
        findings = compliance_check(answer, evidence, kb, verdict)
        blocked = [f for f in findings if f.severity == "block"]
        blocked_total += len(blocked)

        sev_label = {"block": "拦截", "warn": "提示", "info": "说明"}
        if findings:
            finding_text = "\n".join(f"- [{sev_label[f.severity]}] {f.message}" for f in findings)
        else:
            finding_text = "- 未发现明显偏离课件的问题。"
        gate = "**未通过（存在拦截项，请勿直接作为标准答案发布）**" if blocked else "通过"

        banner = provenance_banner(verdict)
        question_recap = f"## 题目复述\n{question}"
        sections.append(
            f"# 第 {idx} 题\n\n{question_recap}\n\n{banner}\n\n{answer}\n\n"
            f"## Harness 合规校验\n- 合规闸门：{gate}\n{finding_text}\n"
        )

    out = Path(args.out).resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text("\n\n".join(sections), encoding="utf-8-sig")
    print(f"Wrote answers: {out}")
    if blocked_total:
        print(f"WARNING: {blocked_total} answer(s) hit a blocking compliance finding; review before publishing.", file=sys.stderr)
    return 0


def command_inspect(args: argparse.Namespace) -> int:
    kb = load_kb(Path(args.kb).resolve())
    print(f"Materials: {len(kb.materials)}")
    print(f"Blocks: {len(kb.blocks)}")
    print(f"Formulas: {len(kb.formulas)}")
    print(f"Format rules: {len(kb.format_rules)}")
    print("\nTop formulas:")
    for formula in kb.formulas[: args.limit]:
        print(f"- {formula}")
    print("\nTop format rules:")
    for rule in kb.format_rules[: args.limit]:
        print(f"- {rule}")
    return 0


def make_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Courseware-grounded standard answer harness for teachers.")
    subparsers = parser.add_subparsers(required=True)

    build = subparsers.add_parser("build", help="Build a courseware knowledge base from PDF/DOCX/PPTX/TXT/MD files.")
    build.add_argument("--materials", nargs="+", required=True, help="Files or folders containing course materials.")
    build.add_argument("--out", default="standard_answer_harness/out/course_kb.json")
    build.add_argument("--render-page-evidence", action="store_true", help="Render PDF pages into image evidence. Slower; use for scanned PDFs after diagnosis.")
    build.add_argument("--page-evidence-dir", help="Where to store rendered page evidence when --render-page-evidence is set.")
    build.set_defaults(func=command_build)

    answer = subparsers.add_parser("answer", help="Generate standard answers from a built knowledge base.")
    answer.add_argument("--kb", required=True)
    answer.add_argument("--questions", required=True)
    answer.add_argument("--out", default="standard_answer_harness/out/answers.md")
    answer.add_argument("--top-k", type=int, default=6)
    answer.add_argument("--model", help="OpenAI model name. If omitted, uses STANDARD_ANSWER_MODEL. If neither is set, runs offline retrieval mode.")
    answer.add_argument("--external-context", help="Optional vetted expansion material. It is treated as secondary to courseware.")
    answer.set_defaults(func=command_answer)

    inspect = subparsers.add_parser("inspect", help="Inspect extracted formulas and format rules.")
    inspect.add_argument("--kb", required=True)
    inspect.add_argument("--limit", type=int, default=20)
    inspect.set_defaults(func=command_inspect)

    diagnose = subparsers.add_parser("diagnose", help="Report a KB's extraction health and whether it can ground answers.")
    diagnose.add_argument("--kb", required=True)
    diagnose.set_defaults(func=command_diagnose)

    diagnose_pdf = subparsers.add_parser("diagnose-pdf", help="Diagnose whether a PDF has usable text or is scanned/image-heavy.")
    diagnose_pdf.add_argument("--pdf", required=True)
    diagnose_pdf.set_defaults(func=command_diagnose_pdf)

    prep_pdf = subparsers.add_parser("prep-pdf", help="Prepare a PDF workspace: safe copy, last-page renders, and contact sheets.")
    prep_pdf.add_argument("--pdf", required=True)
    prep_pdf.add_argument("--out", default="tmp/pdfs/prep")
    prep_pdf.add_argument("--last-pages", type=int, default=3)
    prep_pdf.add_argument("--pages-per-sheet", type=int, default=12)
    prep_pdf.add_argument("--zoom", type=float, default=2.0)
    prep_pdf.set_defaults(func=command_prep_pdf)

    homework = subparsers.add_parser("homework-last-page", help="Render and inspect the last page(s), where homework is often assigned.")
    homework.add_argument("--pdf", required=True)
    homework.add_argument("--out", default="tmp/pdfs/homework_last_page")
    homework.add_argument("--last-pages", type=int, default=1)
    homework.add_argument("--zoom", type=float, default=2.0)
    homework.add_argument("--ocr", action="store_true", help="Run local OCR on rendered last pages if Tesseract is installed.")
    homework.add_argument("--ocr-lang", default="chi_sim+eng")
    homework.set_defaults(func=command_homework_last_page)

    ocr_pages = subparsers.add_parser("ocr-pages", help="Render and OCR selected pages only; never full-document OCR by default.")
    ocr_pages.add_argument("--pdf", required=True)
    ocr_pages.add_argument("--pages", required=True, help="Comma/range list, e.g. 31,32,73-80.")
    ocr_pages.add_argument("--out", default="tmp/pdfs/ocr_pages")
    ocr_pages.add_argument("--zoom", type=float, default=2.0)
    ocr_pages.add_argument("--ocr-lang", default="chi_sim+eng")
    ocr_pages.set_defaults(func=command_ocr_pages)

    locate = subparsers.add_parser("locate-pages", help="Locate likely key pages from the PDF text layer; falls back honestly if no text exists.")
    locate.add_argument("--pdf", required=True)
    locate.add_argument("--query", required=True)
    locate.add_argument("--max-pages", type=int, default=20)
    locate.set_defaults(func=command_locate_pages)

    return parser


def _force_utf8_console() -> None:
    """Windows consoles default to a legacy codepage (e.g. GBK), which turns
    every Chinese diagnostic line into mojibake. Force UTF-8 so the honesty
    warnings are actually readable where they matter most."""
    for stream in (sys.stdout, sys.stderr):
        reconfigure = getattr(stream, "reconfigure", None)
        if reconfigure is not None:
            try:
                reconfigure(encoding="utf-8")
            except (ValueError, OSError):
                pass


def main() -> int:
    _force_utf8_console()
    parser = make_parser()
    args = parser.parse_args()
    return args.func(args)


if __name__ == "__main__":
    raise SystemExit(main())
